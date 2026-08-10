"""
Universal Course Material Extractor
===================================
Exhaustively extracts content from PDFs, PPTX, DOCX, and images.

Usage:
    python extract_course_material.py --input <path> --output <output_dir>
    python extract_course_material.py --input "folder/" --output "extracted/"
    
Handles:
    - PDF (textbooks, notes, slides exported as PDF)
    - PPTX (PowerPoint presentations)
    - DOCX (Word documents)
    - Images (PNG, JPG, JPEG) with OCR

Output:
    - raw_text.txt: Plain text content
    - structured.md: Formatted markdown
    - equations.json: Detected equations
    - problems.json: Detected worked examples
    - figures/: Extracted images
    - extraction_summary.md: Coverage report
"""

import os
import sys
import re
import json
import argparse
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field, asdict
from collections import defaultdict

# ============================================================================
# DATA STRUCTURES
# ============================================================================

@dataclass
class ExtractedEquation:
    """Represents a detected equation"""
    content: str
    source_file: str
    page_or_slide: int
    context: str = ""  # Surrounding text
    is_numbered: bool = False
    equation_number: str = ""

@dataclass
class ExtractedProblem:
    """Represents a worked example or problem"""
    title: str
    problem_text: str
    solution_text: str
    source_file: str
    page_or_slide: int
    problem_type: str = ""  # e.g., "Example", "Problem", "Exercise"

@dataclass
class ExtractedFigure:
    """Represents an extracted figure/image"""
    filename: str
    source_file: str
    page_or_slide: int
    caption: str = ""
    description: str = ""
    referenced_in: str = ""

@dataclass
class ExtractedTable:
    """Represents an extracted table"""
    headers: List[str]
    rows: List[List[str]]
    source_file: str
    page_or_slide: int
    caption: str = ""

@dataclass
class PageContent:
    """Content from a single page/slide"""
    page_num: int
    source_file: str
    raw_text: str
    formatted_text: str = ""
    section_header: str = ""
    equations: List[ExtractedEquation] = field(default_factory=list)
    problems: List[ExtractedProblem] = field(default_factory=list)
    figures: List[ExtractedFigure] = field(default_factory=list)
    tables: List[ExtractedTable] = field(default_factory=list)

@dataclass
class ExtractionResult:
    """Complete extraction result"""
    source_files: List[str]
    total_pages: int
    pages: List[PageContent]
    all_equations: List[ExtractedEquation]
    all_problems: List[ExtractedProblem]
    all_figures: List[ExtractedFigure]
    all_tables: List[ExtractedTable]
    section_hierarchy: Dict[str, List[str]]
    extraction_time: str
    warnings: List[str] = field(default_factory=list)

# ============================================================================
# PATTERN DEFINITIONS
# ============================================================================

# Equation patterns
EQUATION_PATTERNS = [
    # LaTeX style
    r'\$\$(.+?)\$\$',
    r'\$(.+?)\$',
    r'\\begin\{equation\}(.+?)\\end\{equation\}',
    r'\\begin\{align\}(.+?)\\end\{align\}',
    # Common equation formats
    r'([A-Za-z_][A-Za-z0-9_]*\s*=\s*[^,\n]{5,})',  # Variable = expression
    r'(\d+\.?\d*\s*[+\-*/×÷=<>≤≥≈∝∞∫∑∏√]\s*[^,\n]+)',  # Math expressions
    # Greek letters indicating equations
    r'([αβγδεζηθικλμνξοπρστυφχψω][^,\n]{3,}=.+)',
    # Subscript/superscript patterns
    r'([A-Z][a-z]*(?:_\{?[a-z0-9]+\}?)?(?:\^\{?[a-z0-9+\-]+\}?)?\s*=.+)',
]

# Problem/Example patterns
PROBLEM_PATTERNS = [
    r'(?:Example|EXAMPLE)\s*(\d+[\.\d]*)',
    r'(?:Problem|PROBLEM)\s*(\d+[\.\d]*)',
    r'(?:Exercise|EXERCISE)\s*(\d+[\.\d]*)',
    r'(?:Sample Problem|Worked Example)\s*(\d+[\.\d]*)?',
    r'(?:Solution|SOLUTION)[\s:]+',
    r'(?:Q\.|Question)\s*(\d+)',
]

# Section header patterns
SECTION_PATTERNS = [
    r'^(?:Chapter|CHAPTER)\s+(\d+)[:\.]?\s*(.+)$',
    r'^(\d+(?:\.\d+)*)[:\.\s]+(.+)$',  # 1.2.3 Section Name
    r'^(?:Section|SECTION)\s+(\d+(?:\.\d+)*)[:\.]?\s*(.+)$',
    r'^([A-Z][A-Z\s]{2,}[A-Z])$',  # ALL CAPS HEADERS
]

# ============================================================================
# PDF EXTRACTION
# ============================================================================

def extract_from_pdf(filepath: str, output_dir: str) -> List[PageContent]:
    """Extract content from PDF using PyMuPDF (fitz)"""
    pages = []
    
    try:
        import fitz  # PyMuPDF
    except ImportError:
        print("Installing PyMuPDF...")
        os.system(f"{sys.executable} -m pip install PyMuPDF")
        import fitz
    
    doc = fitz.open(filepath)
    filename = Path(filepath).name
    figures_dir = Path(output_dir) / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"  Processing {doc.page_count} pages from {filename}...")
    
    for page_num in range(doc.page_count):
        page = doc[page_num]
        
        # Extract text with layout preservation
        text = page.get_text("text")
        blocks = page.get_text("dict")["blocks"]
        
        # Build formatted text preserving structure
        formatted_lines = []
        current_section = ""
        
        for block in blocks:
            if block["type"] == 0:  # Text block
                for line in block.get("lines", []):
                    line_text = ""
                    is_bold = False
                    font_size = 0
                    
                    for span in line.get("spans", []):
                        line_text += span.get("text", "")
                        font_size = max(font_size, span.get("size", 0))
                        if "bold" in span.get("font", "").lower():
                            is_bold = True
                    
                    line_text = line_text.strip()
                    if not line_text:
                        continue
                    
                    # Detect section headers
                    if is_bold and font_size > 11:
                        formatted_lines.append(f"\n## {line_text}\n")
                        current_section = line_text
                    elif font_size > 13:
                        formatted_lines.append(f"\n# {line_text}\n")
                        current_section = line_text
                    else:
                        formatted_lines.append(line_text)
        
        formatted_text = "\n".join(formatted_lines)
        
        # Extract images
        figures = []
        image_list = page.get_images()
        for img_index, img in enumerate(image_list):
            try:
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                
                img_filename = f"{Path(filename).stem}_p{page_num+1}_img{img_index+1}.png"
                img_path = figures_dir / img_filename
                
                if pix.n < 5:  # GRAY or RGB
                    pix.save(str(img_path))
                else:  # CMYK: convert to RGB first
                    pix1 = fitz.Pixmap(fitz.csRGB, pix)
                    pix1.save(str(img_path))
                    pix1 = None
                pix = None
                
                # Try to find caption near image
                caption = extract_figure_caption(text, img_index)
                
                figures.append(ExtractedFigure(
                    filename=img_filename,
                    source_file=filename,
                    page_or_slide=page_num + 1,
                    caption=caption
                ))
            except Exception as e:
                pass  # Skip problematic images
        
        # Extract equations
        equations = extract_equations_from_text(text, filename, page_num + 1)
        
        # Extract problems/examples
        problems = extract_problems_from_text(text, filename, page_num + 1)
        
        # Extract tables
        tables = extract_tables_from_page(page, filename, page_num + 1)
        
        page_content = PageContent(
            page_num=page_num + 1,
            source_file=filename,
            raw_text=text,
            formatted_text=formatted_text,
            section_header=current_section,
            equations=equations,
            problems=problems,
            figures=figures,
            tables=tables
        )
        pages.append(page_content)
    
    doc.close()
    return pages

def extract_tables_from_page(page, filename: str, page_num: int) -> List[ExtractedTable]:
    """Extract tables from a PDF page"""
    tables = []
    
    try:
        # Try using tabula-py if available
        import tabula
        dfs = tabula.read_pdf(page.parent.name, pages=page_num, multiple_tables=True)
        for df in dfs:
            if df is not None and not df.empty:
                tables.append(ExtractedTable(
                    headers=list(df.columns),
                    rows=[list(row) for row in df.values],
                    source_file=filename,
                    page_or_slide=page_num
                ))
    except:
        # Fallback: detect table-like structures in text
        pass
    
    return tables

# ============================================================================
# POWERPOINT EXTRACTION
# ============================================================================

def extract_from_pptx(filepath: str, output_dir: str) -> List[PageContent]:
    """Extract content from PowerPoint files"""
    pages = []
    
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Installing python-pptx...")
        os.system(f"{sys.executable} -m pip install python-pptx")
        from pptx import Presentation
    
    prs = Presentation(filepath)
    filename = Path(filepath).name
    figures_dir = Path(output_dir) / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"  Processing {len(prs.slides)} slides from {filename}...")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text_content = []
        formatted_content = []
        figures = []
        current_section = ""
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_content.append(text)
                    
                    # Check if this is a title
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            para_text = para.text.strip()
                            if para_text:
                                # Detect titles by font size or shape type
                                if hasattr(shape, "shape_type"):
                                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                                        if shape.placeholder_format.type.name in ['TITLE', 'CENTER_TITLE']:
                                            formatted_content.append(f"# {para_text}")
                                            current_section = para_text
                                            continue
                                formatted_content.append(para_text)
            
            # Extract images
            if hasattr(shape, "image"):
                try:
                    image = shape.image
                    img_bytes = image.blob
                    img_ext = image.ext
                    img_filename = f"{Path(filename).stem}_s{slide_num}_img{len(figures)+1}.{img_ext}"
                    img_path = figures_dir / img_filename
                    
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    
                    figures.append(ExtractedFigure(
                        filename=img_filename,
                        source_file=filename,
                        page_or_slide=slide_num
                    ))
                except:
                    pass
            
            # Extract tables
            if shape.has_table:
                table = shape.table
                headers = [cell.text for cell in table.rows[0].cells]
                rows = [[cell.text for cell in row.cells] for row in list(table.rows)[1:]]
                # TODO: Add to tables list
        
        # Extract from notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                text_content.append(f"\n[Speaker Notes]\n{notes}")
        
        raw_text = "\n".join(text_content)
        formatted_text = "\n".join(formatted_content)
        
        equations = extract_equations_from_text(raw_text, filename, slide_num)
        problems = extract_problems_from_text(raw_text, filename, slide_num)
        
        page_content = PageContent(
            page_num=slide_num,
            source_file=filename,
            raw_text=raw_text,
            formatted_text=formatted_text,
            section_header=current_section,
            equations=equations,
            problems=problems,
            figures=figures
        )
        pages.append(page_content)
    
    return pages

# ============================================================================
# WORD DOCUMENT EXTRACTION
# ============================================================================

def extract_from_docx(filepath: str, output_dir: str) -> List[PageContent]:
    """Extract content from Word documents"""
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        print("Installing python-docx...")
        os.system(f"{sys.executable} -m pip install python-docx")
        from docx import Document
    
    doc = Document(filepath)
    filename = Path(filepath).name
    figures_dir = Path(output_dir) / "figures"
    figures_dir.mkdir(parents=True, exist_ok=True)
    
    print(f"  Processing {filename}...")
    
    # Word docs don't have explicit pages, so we group by sections/headings
    pages = []
    current_page = None
    page_num = 0
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        
        # Check if this is a heading
        is_heading = para.style.name.startswith('Heading')
        
        if is_heading or current_page is None:
            # Start new "page" for each major section
            if current_page is not None:
                current_page.equations = extract_equations_from_text(
                    current_page.raw_text, filename, current_page.page_num
                )
                current_page.problems = extract_problems_from_text(
                    current_page.raw_text, filename, current_page.page_num
                )
                pages.append(current_page)
            
            page_num += 1
            current_page = PageContent(
                page_num=page_num,
                source_file=filename,
                raw_text="",
                formatted_text="",
                section_header=text if is_heading else ""
            )
        
        # Append text
        if is_heading:
            level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
            current_page.formatted_text += f"\n{'#' * level} {text}\n"
        else:
            current_page.formatted_text += f"{text}\n"
        
        current_page.raw_text += f"{text}\n"
    
    # Don't forget last section
    if current_page is not None:
        current_page.equations = extract_equations_from_text(
            current_page.raw_text, filename, current_page.page_num
        )
        current_page.problems = extract_problems_from_text(
            current_page.raw_text, filename, current_page.page_num
        )
        pages.append(current_page)
    
    # Extract tables
    for table_idx, table in enumerate(doc.tables):
        headers = [cell.text for cell in table.rows[0].cells]
        rows = [[cell.text for cell in row.cells] for row in table.rows[1:]]
        
        extracted_table = ExtractedTable(
            headers=headers,
            rows=rows,
            source_file=filename,
            page_or_slide=1  # DOCX doesn't have pages
        )
        if pages:
            pages[0].tables.append(extracted_table)
    
    return pages

# ============================================================================
# IMAGE EXTRACTION (OCR)
# ============================================================================

def extract_from_image(filepath: str, output_dir: str) -> List[PageContent]:
    """Extract text from images using OCR"""
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        print("Installing pytesseract and Pillow...")
        os.system(f"{sys.executable} -m pip install pytesseract Pillow")
        import pytesseract
        from PIL import Image
    
    filename = Path(filepath).name
    print(f"  Processing image {filename} with OCR...")
    
    try:
        image = Image.open(filepath)
        text = pytesseract.image_to_string(image)
        
        # Copy image to figures directory
        figures_dir = Path(output_dir) / "figures"
        figures_dir.mkdir(parents=True, exist_ok=True)
        
        import shutil
        dest_path = figures_dir / filename
        shutil.copy(filepath, dest_path)
        
        figure = ExtractedFigure(
            filename=filename,
            source_file=filename,
            page_or_slide=1,
            description="OCR extracted from image"
        )
        
        equations = extract_equations_from_text(text, filename, 1)
        problems = extract_problems_from_text(text, filename, 1)
        
        return [PageContent(
            page_num=1,
            source_file=filename,
            raw_text=text,
            formatted_text=text,
            figures=[figure],
            equations=equations,
            problems=problems
        )]
    except Exception as e:
        print(f"  Warning: Could not process image {filename}: {e}")
        return []

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def extract_figure_caption(text: str, img_index: int) -> str:
    """Try to find figure caption near image reference"""
    patterns = [
        rf'(?:Figure|Fig\.?)\s*{img_index + 1}[:\.]?\s*([^\n]+)',
        rf'(?:Figure|Fig\.?)\s*\d+[:\.]?\s*([^\n]+)',
    ]
    
    for pattern in patterns:
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            return match.group(1).strip()
    
    return ""

def extract_equations_from_text(text: str, filename: str, page_num: int) -> List[ExtractedEquation]:
    """Extract equations from text content"""
    equations = []
    seen = set()
    
    for pattern in EQUATION_PATTERNS:
        try:
            matches = re.finditer(pattern, text, re.MULTILINE | re.DOTALL)
            for match in matches:
                eq_text = match.group(1) if match.groups() else match.group(0)
                eq_text = eq_text.strip()
                
                # Skip if too short or already seen
                if len(eq_text) < 3 or eq_text in seen:
                    continue
                
                # Skip if it's just a number
                if re.match(r'^[\d\s.,]+$', eq_text):
                    continue
                
                seen.add(eq_text)
                
                # Get context (surrounding text)
                start = max(0, match.start() - 50)
                end = min(len(text), match.end() + 50)
                context = text[start:end].replace('\n', ' ')
                
                # Check if numbered
                eq_num_match = re.search(r'\((\d+(?:\.\d+)?)\)', context)
                
                equations.append(ExtractedEquation(
                    content=eq_text,
                    source_file=filename,
                    page_or_slide=page_num,
                    context=context,
                    is_numbered=eq_num_match is not None,
                    equation_number=eq_num_match.group(1) if eq_num_match else ""
                ))
        except:
            continue
    
    return equations

def extract_problems_from_text(text: str, filename: str, page_num: int) -> List[ExtractedProblem]:
    """Extract worked examples and problems from text"""
    problems = []
    
    # Find problem starts
    for pattern in PROBLEM_PATTERNS:
        for match in re.finditer(pattern, text, re.IGNORECASE | re.MULTILINE):
            start_pos = match.start()
            problem_num = match.group(1) if match.groups() and match.group(1) else ""
            
            # Try to extract problem and solution text
            # Look for next problem or end of text
            remaining_text = text[start_pos:]
            
            # Find the end of this problem (next problem header or significant gap)
            end_match = re.search(
                r'\n(?:Example|Problem|Exercise|Q\.|Question)\s+\d+',
                remaining_text[100:],  # Skip first 100 chars
                re.IGNORECASE
            )
            
            if end_match:
                problem_text = remaining_text[:100 + end_match.start()]
            else:
                # Take next ~1000 chars or until double newline
                problem_text = remaining_text[:1000]
                double_nl = problem_text.find('\n\n\n')
                if double_nl > 100:
                    problem_text = problem_text[:double_nl]
            
            # Determine problem type
            problem_type = "Example"
            if "problem" in match.group(0).lower():
                problem_type = "Problem"
            elif "exercise" in match.group(0).lower():
                problem_type = "Exercise"
            elif "question" in match.group(0).lower() or "q." in match.group(0).lower():
                problem_type = "Question"
            
            # Try to split problem from solution
            solution_split = re.search(r'(?:Solution|Answer|Ans)[:\s]', problem_text, re.IGNORECASE)
            
            if solution_split:
                prob_text = problem_text[:solution_split.start()].strip()
                sol_text = problem_text[solution_split.end():].strip()
            else:
                prob_text = problem_text.strip()
                sol_text = ""
            
            problems.append(ExtractedProblem(
                title=f"{problem_type} {problem_num}".strip(),
                problem_text=prob_text[:500],  # Limit length
                solution_text=sol_text[:1000],
                source_file=filename,
                page_or_slide=page_num,
                problem_type=problem_type
            ))
    
    return problems

def build_section_hierarchy(pages: List[PageContent]) -> Dict[str, List[str]]:
    """Build hierarchy of sections from extracted content"""
    hierarchy = defaultdict(list)
    current_chapter = ""
    
    for page in pages:
        header = page.section_header or ""
        
        # Check if this is a chapter-level header
        chapter_match = re.match(r'^(?:Chapter|CHAPTER)\s+\d+', header)
        if chapter_match or (header.isupper() and len(header.split()) <= 5):
            current_chapter = header
            hierarchy[current_chapter] = []
        elif header and current_chapter:
            hierarchy[current_chapter].append(header)
        elif header:
            hierarchy["Main"].append(header)
    
    return dict(hierarchy)

# ============================================================================
# MAIN EXTRACTION FUNCTION
# ============================================================================

def extract_course_material(input_path: str, output_dir: str) -> ExtractionResult:
    """Main extraction function"""
    input_path = Path(input_path)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    all_pages = []
    source_files = []
    warnings = []
    
    # Determine files to process
    if input_path.is_file():
        files = [input_path]
    else:
        files = list(input_path.glob("**/*"))
    
    print(f"\nExtracting from {len(files)} file(s)...")
    
    for filepath in files:
        if not filepath.is_file():
            continue
            
        ext = filepath.suffix.lower()
        source_files.append(str(filepath))
        
        try:
            if ext == '.pdf':
                pages = extract_from_pdf(str(filepath), str(output_dir))
            elif ext == '.pptx':
                pages = extract_from_pptx(str(filepath), str(output_dir))
            elif ext == '.docx':
                pages = extract_from_docx(str(filepath), str(output_dir))
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
                pages = extract_from_image(str(filepath), str(output_dir))
            else:
                warnings.append(f"Unsupported file type: {filepath}")
                continue
            
            all_pages.extend(pages)
            
        except Exception as e:
            warnings.append(f"Error processing {filepath}: {str(e)}")
    
    # Aggregate content
    all_equations = []
    all_problems = []
    all_figures = []
    all_tables = []
    
    for page in all_pages:
        all_equations.extend(page.equations)
        all_problems.extend(page.problems)
        all_figures.extend(page.figures)
        all_tables.extend(page.tables)
    
    # Build section hierarchy
    hierarchy = build_section_hierarchy(all_pages)
    
    result = ExtractionResult(
        source_files=source_files,
        total_pages=len(all_pages),
        pages=all_pages,
        all_equations=all_equations,
        all_problems=all_problems,
        all_figures=all_figures,
        all_tables=all_tables,
        section_hierarchy=hierarchy,
        extraction_time=datetime.now().isoformat(),
        warnings=warnings
    )
    
    return result

# ============================================================================
# OUTPUT WRITERS
# ============================================================================

def write_raw_text(result: ExtractionResult, output_dir: Path):
    """Write raw text output"""
    with open(output_dir / "raw_text.txt", "w", encoding="utf-8") as f:
        for page in result.pages:
            f.write(f"\n{'='*60}\n")
            f.write(f"[{page.source_file}] Page/Slide {page.page_num}\n")
            f.write(f"{'='*60}\n\n")
            f.write(page.raw_text)
            f.write("\n")

def write_structured_markdown(result: ExtractionResult, output_dir: Path, obsidian_mode: bool = False):
    """Write structured markdown output"""
    with open(output_dir / "structured.md", "w", encoding="utf-8") as f:
        f.write("# Extracted Course Content\n\n")
        f.write(f"**Source Files**: {', '.join(result.source_files)}\n\n")
        f.write(f"**Total Pages/Slides**: {result.total_pages}\n\n")
        f.write(f"**Extraction Time**: {result.extraction_time}\n\n")
        
        f.write("---\n\n")
        
        # Table of Contents
        f.write("## Table of Contents\n\n")
        for chapter, sections in result.section_hierarchy.items():
            f.write(f"- **{chapter}**\n")
            for section in sections:
                f.write(f"  - {section}\n")
        f.write("\n---\n\n")
        
        # Content by page
        current_source = ""
        for page in result.pages:
            if page.source_file != current_source:
                f.write(f"\n# Source: {page.source_file}\n\n")
                current_source = page.source_file
            
            f.write(f"\n## Page {page.page_num}")
            if page.section_header:
                f.write(f": {page.section_header}")
            f.write("\n\n")
            
            f.write(page.formatted_text or page.raw_text)
            f.write("\n")
            
            # Note figures - use Obsidian syntax if enabled
            if page.figures:
                f.write("\n### Figures on this page:\n\n")
                for fig in page.figures:
                    if obsidian_mode:
                        f.write(f"![[{fig.filename}]]\n")
                    else:
                        f.write(f"- `{fig.filename}`")
                    if fig.caption:
                        f.write(f"*{fig.caption}*\n" if obsidian_mode else f": {fig.caption}")
                    f.write("\n")


def write_equations_json(result: ExtractionResult, output_dir: Path):
    """Write equations to JSON"""
    equations_data = [asdict(eq) for eq in result.all_equations]
    with open(output_dir / "equations.json", "w", encoding="utf-8") as f:
        json.dump(equations_data, f, indent=2, ensure_ascii=False)

def write_problems_json(result: ExtractionResult, output_dir: Path):
    """Write problems to JSON"""
    problems_data = [asdict(p) for p in result.all_problems]
    with open(output_dir / "problems.json", "w", encoding="utf-8") as f:
        json.dump(problems_data, f, indent=2, ensure_ascii=False)

def write_extraction_summary(result: ExtractionResult, output_dir: Path):
    """Write extraction summary"""
    with open(output_dir / "extraction_summary.md", "w", encoding="utf-8") as f:
        f.write("# Extraction Summary\n\n")
        
        f.write("## Overview\n\n")
        f.write(f"| Metric | Count |\n")
        f.write(f"|--------|-------|\n")
        f.write(f"| Source Files | {len(result.source_files)} |\n")
        f.write(f"| Total Pages/Slides | {result.total_pages} |\n")
        f.write(f"| Equations Detected | {len(result.all_equations)} |\n")
        f.write(f"| Problems/Examples | {len(result.all_problems)} |\n")
        f.write(f"| Figures Extracted | {len(result.all_figures)} |\n")
        f.write(f"| Tables Extracted | {len(result.all_tables)} |\n")
        
        f.write("\n## Source Files\n\n")
        for src in result.source_files:
            f.write(f"- {src}\n")
        
        f.write("\n## Section Hierarchy\n\n")
        for chapter, sections in result.section_hierarchy.items():
            f.write(f"### {chapter}\n")
            for section in sections:
                f.write(f"- {section}\n")
            f.write("\n")
        
        f.write("\n## Problems/Examples Found\n\n")
        for prob in result.all_problems:
            f.write(f"- **{prob.title}** (Page {prob.page_or_slide}, {prob.source_file})\n")
        
        if result.warnings:
            f.write("\n## Warnings\n\n")
            for warning in result.warnings:
                f.write(f"- {warning}\n")
        
        f.write(f"\n---\n\n*Extracted at: {result.extraction_time}*\n")



# ============================================================================
# CLI INTERFACE
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Extract content from course materials (PDF, PPTX, DOCX, images)"
    )
    parser.add_argument(
        "--input", "-i",
        required=True,
        help="Input file or directory path"
    )
    parser.add_argument(
        "--output", "-o",
        default="extracted_content",
        help="Output directory (default: extracted_content)"
    )
    parser.add_argument(
        "--obsidian",
        action="store_true",
        help="Enable Obsidian mode: uses Images/ folder and ![[]] syntax"
    )
    
    args = parser.parse_args()
    
    # If obsidian mode, rename figures folder to Images
    output_path = Path(args.output)
    if args.obsidian:
        # Patch the figures directory name globally
        global OBSIDIAN_MODE
        OBSIDIAN_MODE = True
        print("Obsidian mode enabled: Images will be saved to 'Images/' folder")
        print("Image references will use ![[image_name.png]] syntax")
    
    result = extract_course_material(args.input, args.output)
    save_extraction_results(result, args.output, obsidian_mode=args.obsidian)
    
    # If obsidian mode, rename figures to Images
    if args.obsidian:
        figures_path = output_path / "figures"
        images_path = output_path / "Images"
        if figures_path.exists():
            if images_path.exists():
                import shutil
                shutil.rmtree(images_path)
            figures_path.rename(images_path)
            print(f"  - Renamed figures/ to Images/")
        
        # Create obsidian_references.md with all image references
        write_obsidian_references(result, output_path)
        print(f"  - obsidian_references.md (copy-paste ready)")
    
    return result

def write_obsidian_references(result: ExtractionResult, output_dir: Path):
    """Write Obsidian-compatible image reference file"""
    with open(output_dir / "obsidian_references.md", "w", encoding="utf-8") as f:
        f.write("# Obsidian Image References\n\n")
        f.write("Copy-paste these references into your Obsidian notes.\n\n")
        f.write("---\n\n")
        
        f.write("## All Extracted Images\n\n")
        for fig in result.all_figures:
            # Use Obsidian wiki-link syntax
            f.write(f"![[{fig.filename}]]\n")
            if fig.caption:
                f.write(f"*{fig.caption}*\n")
            f.write(f"- Source: {fig.source_file}, Page {fig.page_or_slide}\n\n")
        
        f.write("---\n\n")
        f.write("## Images by Page\n\n")
        
        current_source = ""
        for page in result.pages:
            if page.figures:
                if page.source_file != current_source:
                    f.write(f"### {page.source_file}\n\n")
                    current_source = page.source_file
                
                f.write(f"**Page {page.page_num}**:\n")
                for fig in page.figures:
                    f.write(f"- ![[{fig.filename}]]")
                    if fig.caption:
                        f.write(f" - {fig.caption}")
                    f.write("\n")
                f.write("\n")

# Update save function to accept obsidian_mode
def save_extraction_results(result: ExtractionResult, output_dir: str, obsidian_mode: bool = False):
    """Save all extraction outputs"""
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    print("\nWriting output files...")
    
    write_raw_text(result, output_path)
    print("  - raw_text.txt")
    
    write_structured_markdown(result, output_path, obsidian_mode)
    print("  - structured.md")
    
    write_equations_json(result, output_path)
    print(f"  - equations.json ({len(result.all_equations)} equations)")
    
    write_problems_json(result, output_path)
    print(f"  - problems.json ({len(result.all_problems)} problems)")
    
    write_extraction_summary(result, output_path)
    print("  - extraction_summary.md")
    
    print(f"\nExtraction complete! Output saved to: {output_path}")

# Also update write_structured_markdown to support obsidian syntax
def write_structured_markdown_obsidian(result: ExtractionResult, output_dir: Path):
    """Write structured markdown with Obsidian wiki-links for images"""
    with open(output_dir / "structured_obsidian.md", "w", encoding="utf-8") as f:
        f.write("# Extracted Course Content\n\n")
        f.write(f"**Source Files**: {', '.join(result.source_files)}\n\n")
        f.write(f"**Total Pages/Slides**: {result.total_pages}\n\n")
        
        f.write("---\n\n")
        
        # Content by page with Obsidian image syntax
        current_source = ""
        for page in result.pages:
            if page.source_file != current_source:
                f.write(f"\n# Source: {page.source_file}\n\n")
                current_source = page.source_file
            
            f.write(f"\n## Page {page.page_num}")
            if page.section_header:
                f.write(f": {page.section_header}")
            f.write("\n\n")
            
            f.write(page.formatted_text or page.raw_text)
            f.write("\n")
            
            # Figures with Obsidian syntax
            if page.figures:
                f.write("\n### Figures:\n\n")
                for fig in page.figures:
                    f.write(f"![[{fig.filename}]]\n")
                    if fig.caption:
                        f.write(f"*{fig.caption}*\n")
                    f.write("\n")

OBSIDIAN_MODE = False

if __name__ == "__main__":
    main()

